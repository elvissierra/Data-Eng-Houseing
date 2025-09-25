
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_DATA_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.dml.color import RGBColor
from quip2deck.models import SlidePlan
from quip2deck.utils.files import ensure_parent
from quip2deck.settings import RendererSettings
from pathlib import Path
from typing import List, Tuple
import re
import tempfile
import base64
from urllib.parse import urlparse
from urllib.request import urlretrieve

# --- Theme colors ---
BG_DARK = RGBColor(0, 0, 0)
FG_LIGHT = RGBColor(255, 255, 255)

# --- Keynote style constants ---
KEY_FONT = "Helvetica Neue"
TITLE_SIZE_TITLE = Pt(56)   # title slide title
TITLE_SIZE_CONTENT = Pt(48) # content slide title
SUBTITLE_SIZE = Pt(24)
BODY_SIZE = Pt(30)          # body/bullets
SUB_BULLET_SIZE = Pt(24)


def _set_paragraph_text(p, text: str, bold_key_before_colon: bool = True):
    """Set paragraph text with optional bolding of the part before ':'"""
    # Clear any existing simple text
    try:
        p.text = ""
    except Exception:
        pass
    if bold_key_before_colon and ":" in text:
        key, val = text.split(":", 1)
        r1 = p.add_run(); r1.text = key.strip() + ": "
        r1.font.bold = True; r1.font.name = KEY_FONT; r1.font.size = BODY_SIZE
        r2 = p.add_run(); r2.text = val.strip()
        r2.font.bold = False; r2.font.name = KEY_FONT; r2.font.size = BODY_SIZE
    else:
        r = p.add_run(); r.text = text
        r.font.name = KEY_FONT; r.font.size = BODY_SIZE


def _apply_font_paragraph(p, size):
    for r in getattr(p, "runs", []) or []:
        r.font.name = KEY_FONT
        r.font.size = size
        try:
            r.font.color.rgb = FG_LIGHT
        except Exception:
            pass


# --- Dark theme background and text helpers ---
def _apply_dark_background(slide):
    try:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = BG_DARK
    except Exception:
        pass


def _colorize_textframe(tf, rgb):
    try:
        for p in tf.paragraphs:
            for r in getattr(p, 'runs', []) or []:
                if r.font is not None:
                    r.font.color.rgb = rgb
    except Exception:
        pass


# Textual index helper for chart data
def _add_chart_index_box(slide, items: List[Tuple[str, float]], left, top, width, height, show_pct: bool = True):
    if not items:
        return None
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    total = sum(v for (_l, v) in items)
    for i, (lbl, val) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        pct_txt = f" ({val/total:.0%})" if (show_pct and total > 0) else ""
        p.text = f"{lbl} — {int(val) if float(val).is_integer() else val}{pct_txt}"
        p.level = 0
        try:
            for r in p.runs:
                r.font.size = Pt(12)
                try:
                    r.font.color.rgb = FG_LIGHT
                except Exception:
                    pass
        except Exception:
            pass
    return box


def _resolve_image_path(plan: SlidePlan, src: str) -> Path | None:
    """Return a local filesystem Path for an image src.
    - http/https → download to temp and return the temp file path
    - file://    → strip scheme and return local path
    - data:uri   → decode and persist to temp, return path
    - relative   → resolve against plan.meta["base_dir"] if present
    - absolute   → return as-is
    Returns None on failure.
    """
    try:
        if not src:
            return None
        # Normalize common malformed scheme like https:/foo → https://foo
        if src.startswith("https:/") and not src.startswith("https://"):
            src = "https://" + src[len("https:/"):]
        if src.startswith("http:/") and not src.startswith("http://"):
            src = "http://" + src[len("http:/"):]

        # Data URI handling
        if src.startswith("data:"):
            # e.g., data:image/png;base64,AAAA
            try:
                header, b64 = src.split(",", 1)
                # choose extension from header if present
                ext = ".img"
                if ";base64" in header:
                    med = header.split(";")[0]
                else:
                    med = header
                if "/" in med:
                    ext = "." + med.split("/")[-1]
                tmpdir = Path(tempfile.gettempdir()) / "quip2deck_imgs"
                tmpdir.mkdir(parents=True, exist_ok=True)
                tmp_path = tmpdir / ("data_uri" + ext)
                with open(tmp_path, "wb") as f:
                    f.write(base64.b64decode(b64))
                return tmp_path
            except Exception:
                return None

        parsed = urlparse(src)
        scheme = (parsed.scheme or '').lower()
        # Remote URL → download to temp
        if scheme in ("http", "https"):
            tmpdir = Path(tempfile.gettempdir()) / "quip2deck_imgs"
            tmpdir.mkdir(parents=True, exist_ok=True)
            # create a safe filename
            safe_name = re.sub(r"[^A-Za-z0-9_.-]", "_", (parsed.netloc + parsed.path) or "image")
            if not safe_name:
                safe_name = "image"
            tmp_path = tmpdir / safe_name
            headers = {}
            try:
                headers = (plan.meta or {}).get("img_headers") or {}
            except Exception:
                headers = {}
            # Try requests with headers first; fallback to urlretrieve
            try:
                import requests  # type: ignore
                resp = requests.get(src, headers=headers, timeout=20)
                if resp.status_code == 200:
                    tmp_path.write_bytes(resp.content)
                    return tmp_path
            except Exception:
                pass
            try:
                # Fallback without headers
                urlretrieve(src, str(tmp_path))
                return tmp_path
            except Exception:
                return None
        # file:// URL → local path
        if scheme == "file":
            return Path(parsed.path)
        # No scheme → treat as filesystem path
        p = Path(src)
        if p.is_absolute():
            return p
        base_dir = (plan.meta or {}).get("base_dir") if plan else None
        if base_dir:
            return Path(base_dir) / p
        return p
    except Exception:
        return None


def render_pptx(plan: SlidePlan, out_path: str) -> str:
    """Minimal PPTX renderer: Title & Content layout, bullets left, chart right."""
    ensure_parent(out_path)
    # Load renderer settings (overrides may come from plan.meta["settings_override"])
    S = RendererSettings.from_meta(getattr(plan, "meta", None))
    # Rebind global-style constants so helper functions use the theme dynamically
    global BG_DARK, FG_LIGHT, KEY_FONT, TITLE_SIZE_TITLE, TITLE_SIZE_CONTENT, SUBTITLE_SIZE, BODY_SIZE, SUB_BULLET_SIZE
    BG_DARK = RGBColor(*S.bg_rgb)
    FG_LIGHT = RGBColor(*S.fg_rgb)
    KEY_FONT = S.font_name
    TITLE_SIZE_TITLE = Pt(S.title_title_pt)
    TITLE_SIZE_CONTENT = Pt(S.title_content_pt)
    SUBTITLE_SIZE = Pt(S.subtitle_pt)
    BODY_SIZE = Pt(S.body_pt)
    SUB_BULLET_SIZE = Pt(S.sub_bullet_pt)
    prs = Presentation()
    # Use configured slide geometry (defaults preserve 16:9)
    prs.slide_width = Inches(S.slide_width)
    prs.slide_height = Inches(S.slide_height)

    for s in plan.slides:
        if s.layout == "title":
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            _apply_dark_background(slide)
            if slide.shapes.title:
                slide.shapes.title.text = s.title or plan.meta.get("title", "Deck")
                # slightly larger title
                try:
                    for p in slide.shapes.title.text_frame.paragraphs:
                        for r in p.runs:
                            r.font.name = KEY_FONT
                            r.font.size = TITLE_SIZE_TITLE
                except Exception:
                    pass
                _colorize_textframe(slide.shapes.title.text_frame, FG_LIGHT)
            if len(slide.placeholders) > 1 and s.subtitle:
                slide.placeholders[1].text = s.subtitle
                _colorize_textframe(slide.placeholders[1].text_frame, FG_LIGHT)
            continue

        # Content slide
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title & Content
        _apply_dark_background(slide)
        if slide.shapes.title:
            slide.shapes.title.text = s.title or ""
            try:
                for p in slide.shapes.title.text_frame.paragraphs:
                    for r in p.runs:
                        r.font.name = KEY_FONT
                        r.font.size = TITLE_SIZE_CONTENT
            except Exception:
                pass
            _colorize_textframe(slide.shapes.title.text_frame, FG_LIGHT)

        # Optional subtitle: use the body placeholder's first paragraph prefix
        body = slide.placeholders[1].text_frame
        body.clear()
        if s.subtitle:
            body.text = s.subtitle
            try:
                body.paragraphs[0].font.name = KEY_FONT
                body.paragraphs[0].font.size = SUBTITLE_SIZE
            except Exception:
                pass
            # add a spacer paragraph
            body.add_paragraph().text = ""
            try:
                _colorize_textframe(body, FG_LIGHT)
            except Exception:
                pass

        # Paragraphs then bullets (Helvetica Neue), bold key before ':'
        first = (len(body.paragraphs) == 0) or (body.paragraphs[0].text == "")
        if s.paragraphs:
            for ptxt in s.paragraphs:
                p = body.paragraphs[0] if first else body.add_paragraph(); first = False
                p.level = 0
                _set_paragraph_text(p, ptxt, bold_key_before_colon=False)
                _apply_font_paragraph(p, BODY_SIZE)
        if s.bullets:
            for b in s.bullets:
                p = body.paragraphs[0] if first else body.add_paragraph(); first = False
                p.level = 0
                _set_paragraph_text(p, b, bold_key_before_colon=True)
                _apply_font_paragraph(p, BODY_SIZE)

        # Make sure all body text is white
        _colorize_textframe(body, FG_LIGHT)

        # Optional image (placed in the right column). If no chart, image uses the chart box.
        # Collect images
        img_shapes = []
        imgs = s.images or ([s.image] if s.image else [])
        if imgs:
            resolved = []
            missing = []
            for im in imgs:
                pth = _resolve_image_path(plan, im.path)
                if pth and pth.exists():
                    resolved.append((pth, im.alt or ""))
                else:
                    missing.append(im.alt or "Image")

            slide_w = prs.slide_width
            slide_h = prs.slide_height
            margin = Inches(S.margin)
            box_size = min(Inches(S.chart_box_max), slide_h - Inches(3.0))
            left_box = slide_w - box_size - margin
            top_box  = Inches(S.top_box)
            has_chart = bool(getattr(s, "chart", None) and s.chart and s.chart.data)

            if not has_chart:
                if not resolved:
                    tf = slide.shapes.add_textbox(left_box, top_box, box_size, Inches(0.6)).text_frame
                    tf.text = "; ".join([m + " (not available)" for m in missing]) or "Image (not available)"
                    _colorize_textframe(tf, FG_LIGHT)
                else:
                    n = len(resolved)
                    cols = 2 if n >= 2 else 1
                    rows = (n + cols - 1) // cols
                    cell_w = box_size / cols
                    cell_h = box_size / rows
                    for idx, (img_path, alt) in enumerate(resolved):
                        r, c = divmod(idx, cols)
                        left = left_box + int(cell_w * c)
                        top  = top_box  + int(cell_h * r)
                        pic = slide.shapes.add_picture(str(img_path), left, top)
                        scale = min(cell_w / pic.width, cell_h / pic.height, 1.0)
                        pic.width  = int(pic.width  * scale)
                        pic.height = int(pic.height * scale)
                        pic.left = left + int((cell_w - pic.width) / 2)
                        pic.top  = top  + int((cell_h - pic.height) / 2)
                        img_shapes.append(pic)
            else:
                gap = Inches(S.gap)
                thumb_top = top_box + box_size + gap
                bottom_margin = Inches(S.bottom_margin)
                max_top = slide_h - bottom_margin
                avail_h = max_top - thumb_top
                if avail_h > Inches(0.3) and resolved:
                    cols = 2 if len(resolved) >= 2 else 1
                    rows = (len(resolved) + cols - 1) // cols
                    cell_w = box_size / cols
                    cell_h = min(Inches(S.thumb_max_h), avail_h / max(1, rows))
                    for idx, (img_path, alt) in enumerate(resolved):
                        r, c = divmod(idx, cols)
                        left = left_box + int(cell_w * c)
                        top  = thumb_top + int(cell_h * r)
                        pic = slide.shapes.add_picture(str(img_path), left, top)
                        scale = min(cell_w / pic.width, cell_h / pic.height, 1.0)
                        pic.width  = int(pic.width  * scale)
                        pic.height = int(pic.height * scale)
                        pic.left = left + int((cell_w - pic.width) / 2)
                        pic.top  = top  + int((cell_h - pic.height) / 2)
                        img_shapes.append(pic)
                else:
                    # compact placeholder if none resolved or no room
                    tf = slide.shapes.add_textbox(left_box, thumb_top, box_size, Inches(0.4)).text_frame
                    tf.text = "Image" if not resolved else " + ".join([alt or "Image" for _p, alt in resolved])
                    _colorize_textframe(tf, FG_LIGHT)

            for shp in img_shapes:
                try:
                    shp.line.width = Pt(1)
                    shp.line.color.rgb = FG_LIGHT
                except Exception:
                    pass
        # Optional chart on the right
        if getattr(s, "chart", None) and s.chart and s.chart.data:
            cdata = CategoryChartData()
            cdata.categories = [lbl for (lbl, _v) in s.chart.data]
            cdata.add_series("", [v for (_lbl, v) in s.chart.data])
            # right-side placement tuned to 16:9
            slide_w = prs.slide_width
            slide_h = prs.slide_height
            margin = Inches(S.margin)
            chart_size = min(Inches(S.chart_box_max), slide_h - Inches(3.0))  # safe square
            left = slide_w - chart_size - margin
            if left < Inches(0.5):
                # shrink chart to preserve a 0.5" left margin (prevents off-canvas after Keynote reflow)
                overflow = Inches(0.5) - left
                chart_size = max(Inches(3.2), chart_size - overflow)
                left = slide_w - chart_size - margin
                width = chart_size
                height = chart_size
            top = Inches(S.top_box)
            width = chart_size
            height = chart_size
            ctype = XL_CHART_TYPE.COLUMN_CLUSTERED
            if s.chart.type == "bar":
                ctype = XL_CHART_TYPE.BAR_CLUSTERED
            elif s.chart.type == "line":
                ctype = XL_CHART_TYPE.LINE
            elif s.chart.type == "pie":
                ctype = XL_CHART_TYPE.PIE
            chart = slide.shapes.add_chart(ctype, left, top, width, height, cdata).chart

            try:
                if s.chart.type == "pie":
                    # Show native legend (color index) and keep percentage labels
                    chart.has_legend = True
                    try:
                        pos = {
                            "bottom": XL_LEGEND_POSITION.BOTTOM,
                            "right": XL_LEGEND_POSITION.RIGHT,
                            "left": XL_LEGEND_POSITION.LEFT,
                            "top": XL_LEGEND_POSITION.TOP,
                        }.get("bottom", XL_LEGEND_POSITION.BOTTOM)
                        chart.legend.position = pos
                    except Exception:
                        pass
                    plot = chart.plots[0]
                    plot.has_data_labels = True
                    d = plot.data_labels
                    d.show_percentage = True
                    d.number_format = "0%"
                    try:
                        d.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
                    except Exception:
                        pass
                    # White labels/legend for dark background
                    try:
                        d.font.color.rgb = FG_LIGHT
                    except Exception:
                        pass
                    try:
                        chart.legend.font.color.rgb = FG_LIGHT
                    except Exception:
                        pass
                    # Apply settings overrides (legend on/off, position, label style)
                    try:
                        chart.has_legend = bool(S.pie_show_legend)
                        pos = {
                            "bottom": XL_LEGEND_POSITION.BOTTOM,
                            "right": XL_LEGEND_POSITION.RIGHT,
                            "left": XL_LEGEND_POSITION.LEFT,
                            "top": XL_LEGEND_POSITION.TOP,
                        }.get((S.legend_position or "bottom").lower(), XL_LEGEND_POSITION.BOTTOM)
                        chart.legend.position = pos
                    except Exception:
                        pass
                    try:
                        d.show_percentage = bool(S.show_percentages)
                        if S.show_percentages:
                            d.number_format = "0%"
                        if S.pie_labels_outside:
                            d.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
                    except Exception:
                        pass
                else:
                    for series in chart.series:
                        series.has_data_labels = True
                        series.data_labels.show_value = True
                        series.data_labels.font.size = Pt(12)
                    # Axes/legend tick labels in white
                    try:
                        if chart.has_legend and chart.legend is not None:
                            chart.legend.font.color.rgb = FG_LIGHT
                        if hasattr(chart, 'category_axis') and chart.category_axis is not None:
                            chart.category_axis.tick_labels.font.color.rgb = FG_LIGHT
                        if hasattr(chart, 'value_axis') and chart.value_axis is not None:
                            chart.value_axis.tick_labels.font.color.rgb = FG_LIGHT
                    except Exception:
                        pass
            except Exception:
                pass



    prs.save(out_path)
    return out_path
